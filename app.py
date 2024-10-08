import streamlit as st
from docx import Document
from pptx import Presentation
import io
import base64

st.markdown(
     """
    <style>
    header[data-testid="stHeader"] {
        display: none;
    }
    .css-1u7wa8r.e1tzin5v3 {
        display: none;
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Fetch password from Streamlit secrets
PASSWORD = "imr solution"

# Password input
password = st.sidebar.text_input("Password", type="password")

# Check password
if password == PASSWORD:
    st.title("Document Text Replacer")

    # Options for segments and their subsegments
    segment_options_ordered = ['SEGMENTTA', 'SEGMENTTB', 'SEGMENTTC', 'SEGMENTTD', 'SEGMENTTE', 'SEGMENTTF']
    segment_options = {
        'SEGMENTTA': ['SEGMENTTA','SUBSEGA1', 'SUBSEGA2', 'SUBSEGA3', 'SUBSEGA4', 'SUBSEGA5', 'SUBSEGA6'],
        'SEGMENTTB': ['SEGMENTTB','SUBSEGB1', 'SUBSEGB2', 'SUBSEGB3', 'SUBSEGB4', 'SUBSEGB5', 'SUBSEGB6'],
        'SEGMENTTC': ['SEGMENTTC','SUBSEGC1', 'SUBSEGC2', 'SUBSEGC3', 'SUBSEGC4', 'SUBSEGC5', 'SUBSEGC6'],
        'SEGMENTTD': ['SEGMENTTD','SUBSEGD1', 'SUBSEGD2', 'SUBSEGD3', 'SUBSEGD4', 'SUBSEGD5', 'SUBSEGD6'],
        'SEGMENTTE': ['SEGMENTTE','SUBSEGE1', 'SUBSEGE2', 'SUBSEGE3', 'SUBSEGE4', 'SUBSEGE5', 'SUBSEGE6'],
        'SEGMENTTF': ['SEGMENTTF','SUBSEGF1', 'SUBSEGF2', 'SUBSEGF3', 'SUBSEGF4', 'SUBSEGF5', 'SUBSEGF6'],
    }
    company_options = ['COMPANYA', 'COMPANYB', 'COMPANYC', 'COMPANYD', 'COMPANYE', 'COMPANYF', 'COMPANYG', 'COMPANYH', 'COMPANYI', 'COMPANYJ', 'COMPANYK', 'COMPANYL', 'COMPANYM', 'COMPANYN', 'COMPANYO', 'COMPANYP', 'COMPANYQ', 'COMPANYR', 'COMPANYS', 'COMPANYT']

    def get_segments_up_to(selected_segment):
        """Returns a list of all segments up to the selected segment."""
        selected_index = segment_options_ordered.index(selected_segment)
        return segment_options_ordered[:selected_index + 1]

    def replace_text_case_sensitive(paragraphs, find_str, replace_str):
        """Replace occurrences of find_str with replace_str in a case-sensitive manner within paragraphs."""
        for para in paragraphs:
            # Combine all runs in a paragraph into a single string
            combined_text = "".join(run.text for run in para.runs)

            # Case-sensitive replacement
            if find_str in combined_text:
                # Replace the found text with the replacement string
                combined_text = combined_text.replace(find_str, replace_str)

                # Clear the existing runs
                para.clear()

                # Re-add the replaced text, preserving formatting
                for part in combined_text.splitlines(True):
                    run = para.add_run(part)
                    run.font.name = "Segoe UI"

    def replace_word_in_docx(doc, find_replace_pairs):
        """
        Perform find and replace in a Word document, maintaining case-sensitive replacements.
        """
        for find_str, replace_str in find_replace_pairs:
            # Replace in regular paragraphs
            replace_text_case_sensitive(doc.paragraphs, find_str, replace_str)
            
            # Replace in table cells
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        replace_text_case_sensitive(cell.paragraphs, find_str, replace_str)
                        
            # Replace in headers and footers
            for section in doc.sections:
                replace_text_case_sensitive(section.header.paragraphs, find_str, replace_str)
                replace_text_case_sensitive(section.footer.paragraphs, find_str, replace_str)

    def replace_text_in_pptx(slides, find_replace_pairs):
        """
        Perform a case-sensitive find and replace operation in PowerPoint slides.
        Only exact case matches are replaced.
        """
        for slide in slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    for find_str, replace_str in find_replace_pairs:
                        # Case-sensitive replacement
                        if find_str in text:
                            paragraph.text = text.replace(find_str, replace_str)

    st.sidebar.header("Upload File")
    uploaded_file = st.sidebar.file_uploader("Upload a File", type=["docx", "pptx"])
    
    st.sidebar.header("Find and Replace")
    user_find = st.sidebar.text_input("Find:")
    user_replace = st.sidebar.text_input("Replace with:")
    
    st.sidebar.header("Second Find and Replace")
    second_find = st.sidebar.text_input("Find (Second):")
    second_replace = st.sidebar.text_input("Replace with (Second):")
    
    st.sidebar.header("Select Segments")
    selected_segment = st.sidebar.selectbox("Select a segment", options=segment_options_ordered)
    
    if selected_segment:
        selected_segments = get_segments_up_to(selected_segment)
        st.sidebar.write(f"Selected Segments: {', '.join(selected_segments)}")
    
    with st.expander("Segment Replacements"):
        segment_replace_inputs = {}
        for segment in selected_segments:
            st.subheader(segment)
            segment_replace_inputs[segment] = {}
            for subsegment in segment_options[segment]:
                segment_replace_inputs[segment][subsegment] = st.text_input(f"Replace {subsegment} with:", key=f"{segment}_{subsegment}")
    
    with st.expander("Company Replacements"):
        company_replace_inputs = {value: st.text_input(f"Replace {value} with:", key=f"COMPANY_{value}") for value in company_options}
    
    st.sidebar.header("Download Options")
    custom_filename = st.sidebar.text_input("Enter Filename:", "")
    
    if st.button("Update File"):
        if uploaded_file:
            try:
                file_content = io.BytesIO(uploaded_file.getvalue())
                filename = uploaded_file.name
                
                # Gather all find-replace pairs, including both sets of inputs
                find_replace_pairs = [(find_str, replace_str) for segment in segment_replace_inputs.values() for find_str, replace_str in segment.items() if replace_str]
                
                for find_str, replace_str in company_replace_inputs.items():
                    if replace_str:
                        find_replace_pairs.append((find_str, replace_str))
                
                if user_find and user_replace:
                    find_replace_pairs.append((user_find, user_replace))
                
                if second_find and second_replace:
                    find_replace_pairs.append((second_find, second_replace))
                
                if filename.endswith('.docx'):
                    doc = Document(file_content)
                    replace_word_in_docx(doc, find_replace_pairs)
                    output_buffer = io.BytesIO()
                    doc.save(output_buffer)
                    output_buffer.seek(0)
                    
                    b64 = base64.b64encode(output_buffer.read()).decode()
                    href = f'<a download="{custom_filename if custom_filename else "modified_document"}.docx" href="data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{b64}" target="_blank">Download Updated Word File</a>'
                    st.markdown(href, unsafe_allow_html=True)
                
                elif filename.endswith('.pptx'):
                    ppt = Presentation(file_content)
                    replace_text_in_pptx(ppt.slides, find_replace_pairs)
                    output_buffer = io.BytesIO()
                    ppt.save(output_buffer)
                    output_buffer.seek(0)
                    
                    b64 = base64.b64encode(output_buffer.read()).decode()
                    href = f'<a download="{custom_filename if custom_filename else "modified_presentation"}.pptx" href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" target="_blank">Download Updated PowerPoint File</a>'
                    st.markdown(href, unsafe_allow_html=True)
                
            except Exception as e:
                st.error(f"An error occurred: {e}")
        else:
            st.error("Upload a File.")
else:
    st.warning("Please enter the correct password.")
